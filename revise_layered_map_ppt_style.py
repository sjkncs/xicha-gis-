# -*- coding: utf-8 -*-
"""Restyle the layered GIS map slides to match the HIT defense PPT template.

The previous insertion used a new rounded side panel and a grey footer, which
made slides 15-20 visually inconsistent with the surrounding deck. This script
rebuilds only those six slides in-place using the same header/footer geometry
and colors as slide 14, while adding concrete spatial interpretation.
"""

from __future__ import annotations

import io
import shutil
from datetime import datetime
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
FIG_DIR = ROOT / "projects" / "15min-urban-accessibility" / "paper" / "figures" / "layered_total_map"

TEAL = RGBColor(0x00, 0x53, 0x75)
ORANGE = RGBColor(0xE6, 0x8A, 0x21)
NAVY = RGBColor(0x1B, 0x2C, 0x35)
TEXT = RGBColor(0x1B, 0x2C, 0x35)
MUTED = RGBColor(0x5F, 0x6F, 0x78)
GREY = RGBColor(0x8A, 0x96, 0x9C)
BORDER = RGBColor(0xE6, 0xEC, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_ORANGE = RGBColor(0xFE, 0xF3, 0xE7)


def latest_target_ppt() -> Path:
    candidates = [
        p
        for p in ROOT.glob("*.pptx")
        if ".before_" not in p.name and not p.name.endswith(".tmp.pptx")
    ]
    if not candidates:
        raise FileNotFoundError("No editable PPTX found in current directory.")
    return max(candidates, key=lambda p: p.stat().st_mtime)


def set_typeface(run, font_name: str) -> None:
    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = r_pr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            r_pr.append(node)
        node.set("typeface", font_name)


def style_run(run, font_name: str, size: float, color: RGBColor, bold: bool = False) -> None:
    set_typeface(run, font_name)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    font_size: float = 14,
    font_name: str = "宋体",
    color: RGBColor = TEXT,
    bold: bool = False,
    align: PP_ALIGN | None = None,
    line_spacing: float = 1.08,
    space_after: float = 2,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)

    lines = text.split("\n") if text else [""]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        if align is not None:
            p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for run in p.runs:
            style_run(run, font_name, font_size, color, bold)
    return box


def add_rect(slide, left: int, top: int, width: int, height: int, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    return shape


def add_line(slide, x1: int, y1: int, x2: int, y2: int, color: RGBColor = TEAL):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(0.9)
    return line


def clear_slide(slide) -> None:
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape._element)


def add_template_header(slide, prs: Presentation, section: str, title: str, page_no: int, logo_blob: bytes | None) -> None:
    sw = prs.slide_width
    sh = prs.slide_height
    add_rect(slide, 0, 0, sw, sh, WHITE)
    add_rect(slide, 0, 0, sw, 91440, TEAL)
    add_text(slide, "≡", 274320, 246888, 182880, 146304, font_size=18, font_name="黑体", color=TEAL, bold=True)
    add_text(slide, section, 676656, 201168, 603504, 182880, font_size=17, font_name="黑体", color=ORANGE, bold=True)
    add_text(slide, title, 1298448, 192024, 7955279, 201168, font_size=16.5, font_name="黑体", color=NAVY, bold=True)
    add_line(slide, 658368, 576072, 658368 + 10881360, 576072)
    if logo_blob:
        slide.shapes.add_picture(io.BytesIO(logo_blob), 10131552, 164592, width=1316736, height=392711)
    add_text(
        slide,
        str(page_no),
        11631168,
        237744,
        274320,
        137160,
        font_size=12,
        font_name="Microsoft YaHei",
        color=RGBColor(0x11, 0x18, 0x27),
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_template_footer(slide, conclusion: str, label: str) -> None:
    add_line(slide, 640080, 6144768, 640080 + 10899648, 6144768)
    add_rect(slide, 749808, 6272784, 960120, 274320, TEAL)
    add_text(slide, "核心结论", 813816, 6345936, 841248, 100584, font_size=9.5, font_name="黑体", color=WHITE, bold=True)
    add_text(slide, conclusion, 1847088, 6272784, 8138160, 246888, font_size=13.1, font_name="宋体", color=TEXT)
    add_text(slide, label, 10040112, 6327648, 1426464, 118872, font_size=9.5, font_name="宋体", color=GREY, align=PP_ALIGN.CENTER)


def add_info_card(
    slide,
    title: str,
    body_lines: list[str],
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    accent: RGBColor = TEAL,
    body_fill: RGBColor = WHITE,
    body_font_size: float = 10.3,
) -> None:
    header_h = 310896
    add_rect(slide, left, top, width, height, body_fill, BORDER)
    add_rect(slide, left, top, width, header_h, accent)
    add_text(slide, title, left + 91440, top + 70400, width - 182880, 150000, font_size=12.8, font_name="黑体", color=WHITE, bold=True)
    body = "\n".join(body_lines)
    add_text(
        slide,
        body,
        left + 91440,
        top + header_h + 91440,
        width - 182880,
        height - header_h - 128016,
        font_size=body_font_size,
        font_name="宋体",
        color=TEXT,
        line_spacing=1.08,
        space_after=3,
    )


def add_cropped_picture_to_fit(
    slide,
    image_path: Path,
    left: int,
    top: int,
    max_width: int,
    max_height: int,
    *,
    crop_top: float = 0.12,
    crop_bottom: float = 0.025,
):
    with Image.open(image_path) as img:
        visible_ratio = img.width / (img.height * (1 - crop_top - crop_bottom))
    slot_ratio = max_width / max_height
    if slot_ratio > visible_ratio:
        height = max_height
        width = int(height * visible_ratio)
    else:
        width = max_width
        height = int(width / visible_ratio)
    pic_left = left + int((max_width - width) / 2)
    pic_top = top + int((max_height - height) / 2)
    pic = slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=width, height=height)
    pic.crop_top = crop_top
    pic.crop_bottom = crop_bottom
    return pic


def find_logo_blob(prs: Presentation) -> bytes | None:
    if len(prs.slides) < 14:
        return None
    for shape in prs.slides[13].shapes:
        if shape.shape_type == 13 and shape.left > 9_000_000 and shape.top < 500_000:
            return shape.image.blob
    return None


SLIDES = [
    {
        "index": 14,
        "page": 15,
        "section": "3.2A",
        "title": "多源空间证据链：分层读图逻辑",
        "headline": "从“总图叠加”改为“图层证据链”，先看结构，再看供给、路径与风险",
        "subtitle": "本研究总图按可达性幻觉形成过程分层说明，避免把全部点线面压在同一图幅内",
        "image": "fig_total_layer_06_four_panel_evidence_chain.png",
        "card1_title": "图幅判读",
        "card1": [
            "四联图依次对应空间骨架、服务供给、路径观测和风险结果。",
            "底层数据含13,911条道路、66,424个POI、24,951条路由边和240,458个街景采样点；若直接全量叠加，路网与POI会遮蔽社区风险。",
            "因此本页先建立“结构—供给—路径—风险”的读图顺序。",
        ],
        "card2_title": "异常线索",
        "card2": [
            "高风险并非随机分布：翠湖苑、花样年苑位于西丽/科技园边缘，企业与交通类POI密集但夜间服务不足。",
            "万科村、龙瑞花园位于前海/桂湾骨架路附近，路由边较多但慢行连续性不足。",
            "这些线索说明，单看POI密度或单看道路数量都会低估真实通行问题。",
        ],
        "conclusion": "分层读图把服务、路径、时间和受影响社区拆开，形成后续指标计算的空间证据链。",
        "label": "图3-2A",
        "crop_top": 0.11,
        "crop_bottom": 0.02,
    },
    {
        "index": 15,
        "page": 16,
        "section": "3.2B",
        "title": "空间基础层：路网等级与建筑/AOI",
        "headline": "道路等级决定可达性上限，建筑/AOI揭示真实阻隔与用地边界",
        "subtitle": "本页解释为什么直线距离会高估服务可达：路网等级、街区孔隙与AOI边界共同影响步行成本",
        "image": "fig_total_layer_01_base_road_hierarchy.png",
        "card1_title": "图幅判读",
        "card1": [
            "路网总计13,911条，其中步行/台阶/绿道5,615条、居住/服务道路4,429条、主干/快速路1,758条。",
            "深色骨架路连接片区，浅色服务路和紫色慢行线才决定社区最后500m的真实可达。",
            "建筑/AOI在科技园、前海、蛇口等片区形成连续边界，解释了地图上“近而绕”的空间形态。",
        ],
        "card2_title": "异常定位",
        "card2": [
            "翠湖苑靠近科苑北路、科慧东路、高新中二道/三道，路由边304条、慢行边184条，但园区边界和主干路切割仍抬高夜间SAII。",
            "龙瑞花园沿怡海大道、前海大道、鲤鱼门街，障碍边57条；卓越名苑在后海大道、登良路附近出现2个断头节点。",
            "这些位置说明“路多”不等于“走得通”。",
        ],
        "conclusion": "空间基础层揭示主干路切割、慢行缺口和 AOI 阻隔，是 Network Ratio 修正的依据。",
        "label": "图3-2B",
        "crop_top": 0.12,
        "crop_bottom": 0.02,
    },
    {
        "index": 16,
        "page": 17,
        "section": "3.2C",
        "title": "服务供给层：POI密度与设施类型结构",
        "headline": "POI密度反映服务聚集，但不能直接等同于可用服务",
        "subtitle": "本页从路网中拆出服务供给，识别“密度高、夜间弱、类别偏企业/交通”的可达性幻觉",
        "image": "fig_total_layer_02_poi_service_supply.png",
        "card1_title": "图幅判读",
        "card1": [
            "POI图层共66,424个点，购物13,081个、餐饮12,576个、交通5,526个、教育3,163个、医疗2,556个、公共/生活1,802个。",
            "南部和中部热度高，北部及山体边缘较稀疏；同一热区内，POI类别差异决定其是否能支撑日常生活服务。",
        ],
        "card2_title": "异常定位",
        "card2": [
            "翠湖苑800m内POI 847个，但夜间可用仅25个，夜间率3.0%；近邻多为深圳太阳生物科技、赛诺生、天道医药等公司企业。",
            "花样年苑POI 1,195个、夜间36个，同为3.0%，周边集中在北环科苑天桥、大族激光、停车场。",
            "卓越名苑POI 3,837个但夜间率7.8%，仍需与后海大道/登良路的路径阻隔共同解释。",
        ],
        "conclusion": "高密度 POI 不能直接等同于可用服务，仍需夜间可用性和步行网络共同校准。",
        "label": "图3-2C",
        "crop_top": 0.12,
        "crop_bottom": 0.02,
    },
    {
        "index": 17,
        "page": 18,
        "section": "3.2D",
        "title": "路径证据层：步行路由与街景采样轨迹",
        "headline": "路径、节点和街景采样共同揭示“看似近、实际绕”的通行异常",
        "subtitle": "本页将步行路由和20m街景采样点作为 Network Ratio、M2SFCA 与 GTA 的空间证据",
        "image": "fig_total_layer_03_walk_route_streetview_sampling.png",
        "card1_title": "图幅判读",
        "card1": [
            "青色路由线与橙色街景采样点不是装饰图层，而是指标计算的观测骨架。",
            "路由线进入Network Ratio/M2SFCA，20m街景点进入GTA与步行环境校准。",
            "图中轨迹沿主干和居住街巷展开，可看出南山南部采样密、北部与边缘地带采样较稀疏。",
        ],
        "card2_title": "异常定位",
        "card2": [
            "万科村500m内路由边172条，但慢行边只有26条，桂湾五路、桂湾四路、月湾西街周边障碍边44条，体现前海骨架路下的慢行断裂。",
            "宝安新区在中山公园外围、南头文化街、河滨路一带障碍边84条，说明公园/旧城边界造成路径碎片化。",
            "卓越名苑500m内有2个断头节点，后海大道—登良路附近存在局部绕行异常。",
        ],
        "conclusion": "路由边、慢行边、障碍边和断头节点共同决定真实步行成本，是网络修正的直接证据。",
        "label": "图3-2D",
        "crop_top": 0.12,
        "crop_bottom": 0.02,
    },
    {
        "index": 18,
        "page": 19,
        "section": "3.2E",
        "title": "风险解释层：社区/AOI与SAII时间贫困",
        "headline": "SAII将服务时间、路径连续性和空间阻隔落到具体社区",
        "subtitle": "本页用社区尺度解释哪些居民承受了可达性幻觉，以及异常来自哪类空间机制",
        "image": "fig_total_layer_04_community_aoi_saii_risk.png",
        "card1_title": "图幅判读",
        "card1": [
            "风险点颜色越红表示SAII越高，圆点大小表示社区人口规模。",
            "全域391个社区中，西丽均值SAII为0.1029，明显高于蛇口0.0195、南山0.0192、科技园0.0187、粤海0.0169和沙河0.0158。",
            "风险并不简单随中心区服务密度升降，而与边缘区服务时间和路网结构有关。",
        ],
        "card2_title": "异常定位",
        "card2": [
            "最高风险点依次为翠湖苑0.1702、万科村0.1682、花样年苑0.1630、龙瑞大厦0.1536、龙瑞花园0.1510。",
            "翠湖苑/花样年苑是“企业、停车、交通POI密集但夜间服务弱”；万科村/龙瑞花园是“前海骨架路多但慢行连接不足”。",
            "宝安新区则表现为公园与旧城边界的障碍边聚集。",
        ],
        "conclusion": "SAII 将服务时间、路径连续性和空间阻隔落到社区，解释高风险点的形成机制。",
        "label": "图3-2E",
        "crop_top": 0.12,
        "crop_bottom": 0.02,
    },
    {
        "index": 19,
        "page": 20,
        "section": "3.2F",
        "title": "低密度综合图：从图层证据进入指标计算",
        "headline": "综合图只保留关键证据，服务于后续公式推导而不是再次堆叠所有对象",
        "subtitle": "本页把供给、路径和风险压缩到一张可读图上，承接下一节网络距离修正",
        "image": "fig_total_layer_05_low_density_synthesis.png",
        "card1_title": "图幅判读",
        "card1": [
            "低密度综合图只保留道路骨架、POI热区、步行路由和高风险社区。",
            "它避免再次把6.6万个POI和24万街景点全叠在一张图上。",
            "因此，本图的作用不是展示所有对象，而是把“供给—路径—风险”的主要空间关系压缩成答辩可读的证据图。",
        ],
        "card2_title": "指标承接",
        "card2": [
            "翠湖苑/花样年苑体现“POI密集但夜间率只有3.0%”，对应M2SFCA夜间供给修正。",
            "万科村慢行边26条、障碍边44条，对应Network Ratio路径修正；卓越名苑断头节点2个、障碍边36条，对应GTA与街景环境校准。",
            "下一页进入网络距离与欧氏距离修正，就是从这张综合图提炼出的第一层公式。",
        ],
        "conclusion": "综合图把异常社区、服务热区和路径证据放在同一逻辑面上，承接指标计算章节。",
        "label": "图3-2F",
        "crop_top": 0.12,
        "crop_bottom": 0.02,
    },
]


def rebuild_slide(slide, prs: Presentation, logo_blob: bytes | None, spec: dict) -> None:
    clear_slide(slide)
    add_template_header(slide, prs, spec["section"], spec["title"], spec["page"], logo_blob)

    add_text(slide, spec["headline"], 786384, 749808, 10241280, 347472, font_size=20, font_name="黑体", color=TEAL, bold=True)
    add_text(slide, spec["subtitle"], 804672, 1143000, 10789920, 182880, font_size=11.2, font_name="宋体", color=MUTED)

    image_path = FIG_DIR / spec["image"]
    if not image_path.exists():
        raise FileNotFoundError(image_path)
    add_cropped_picture_to_fit(
        slide,
        image_path,
        Inches(0.62),
        Inches(1.39),
        Inches(5.85),
        Inches(4.42),
        crop_top=spec.get("crop_top", 0.12),
        crop_bottom=spec.get("crop_bottom", 0.02),
    )

    card_left = Inches(6.65)
    card_width = Inches(5.42)
    add_info_card(
        slide,
        spec["card1_title"],
        spec["card1"],
        card_left,
        Inches(1.39),
        card_width,
        Inches(2.08),
        accent=TEAL,
        body_fill=WHITE,
        body_font_size=10.1,
    )
    add_info_card(
        slide,
        spec["card2_title"],
        spec["card2"],
        card_left,
        Inches(3.67),
        card_width,
        Inches(2.18),
        accent=ORANGE,
        body_fill=SOFT_ORANGE,
        body_font_size=10.0,
    )
    add_template_footer(slide, spec["conclusion"], spec["label"])


def main() -> None:
    ppt = latest_target_ppt()
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup = ppt.with_name(f"{ppt.stem}.before_layered_map_style_fix_{timestamp}{ppt.suffix}")
    shutil.copy2(ppt, backup)

    prs = Presentation(str(ppt))
    if len(prs.slides) < 20:
        raise RuntimeError(f"Expected at least 20 slides, got {len(prs.slides)}")

    logo_blob = find_logo_blob(prs)
    for spec in SLIDES:
        rebuild_slide(prs.slides[spec["index"]], prs, logo_blob, spec)

    prs.save(str(ppt))
    print(f"updated: {ppt}")
    print(f"backup: {backup}")
    print("rebuilt slides: 15-20")


if __name__ == "__main__":
    main()
