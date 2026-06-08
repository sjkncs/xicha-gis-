# -*- coding: utf-8 -*-
"""
Insert layered GIS maps into the current HIT defense PPTX.

Insertion point:
  after existing slide 14 (data source summary), before network-distance method.
"""

from __future__ import annotations

import io
import shutil
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


ROOT = Path(__file__).resolve().parent
FIG_DIR = ROOT / "projects" / "15min-urban-accessibility" / "paper" / "figures" / "layered_total_map"

NAVY = RGBColor(17, 24, 39)
MUTED = RGBColor(55, 65, 81)
LIGHT_BG = RGBColor(248, 250, 252)
ACCENT = RGBColor(31, 78, 121)
GOLD = RGBColor(191, 145, 42)
WHITE = RGBColor(255, 255, 255)


def latest_target_ppt() -> Path:
    candidates = [
        p
        for p in ROOT.glob("*.pptx")
        if ".before_" not in p.name and not p.name.endswith(".tmp.pptx")
    ]
    if not candidates:
        raise FileNotFoundError("No editable PPTX found in current directory.")
    return max(candidates, key=lambda p: p.stat().st_mtime)


def add_textbox(slide, text, left, top, width, height, font_size=18, bold=False, color=MUTED, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    for run in p.runs:
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_multiline(slide, lines, left, top, width, height, font_size=14.5, color=MUTED):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(4)
        for run in p.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return box


def add_header(slide, section, title, page_no, logo_blob=None):
    add_textbox(slide, "≡", 274320, 246888, 182880, 146304, font_size=18, bold=True, color=ACCENT)
    add_textbox(slide, section, 676656, 201168, 603504, 182880, font_size=17, bold=True, color=ACCENT)
    add_textbox(slide, title, 1298448, 192024, 7955279, 201168, font_size=22, bold=True, color=NAVY)
    if logo_blob:
        slide.shapes.add_picture(io.BytesIO(logo_blob), 10131552, 164592, width=1316736, height=392711)
    add_textbox(slide, str(page_no), 11631168, 237744, 274320, 137160, font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def add_side_panel(slide, heading, bullets, conclusion):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 8732520, 1264920, 2834640, 4358640)
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT_BG
    panel.line.color.rgb = RGBColor(222, 226, 232)
    panel.line.width = Pt(1.0)
    add_textbox(slide, heading, 8950000, 1463040, 2390000, 260000, font_size=16, bold=True, color=NAVY)
    add_multiline(slide, bullets, 8950000, 1874520, 2390000, 2460000, font_size=13.2, color=MUTED)
    add_textbox(slide, "承接关系", 8950000, 4495800, 900000, 180000, font_size=12.5, bold=True, color=ACCENT)
    add_textbox(slide, conclusion, 8950000, 4770120, 2390000, 620000, font_size=12.8, bold=True, color=NAVY)


def add_footer(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 6228000, 12191365, 628000)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(244, 247, 250)
    bar.line.fill.background()
    add_textbox(slide, "核心结论", 813816, 6345936, 841248, 100584, font_size=9.5, bold=True, color=ACCENT)
    add_textbox(slide, text, 1847088, 6272784, 8138160, 246888, font_size=14.0, bold=True, color=NAVY)


def add_cropped_picture(slide, image_path, left, top, width, height, crop_top=0.13, crop_bottom=0.03):
    pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    pic.crop_top = crop_top
    pic.crop_bottom = crop_bottom
    return pic


def add_map_slide(prs, logo_blob, section, title, image_name, heading, bullets, conclusion, page_no, crop_top=0.13):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, section, title, page_no, logo_blob)
    add_textbox(
        slide,
        "分层展示真实图层，避免把路网、POI、路径和风险点压在同一张总图中。",
        786384,
        749808,
        10241280,
        210000,
        font_size=13.5,
        color=MUTED,
    )
    add_cropped_picture(
        slide,
        FIG_DIR / image_name,
        640080,
        1264920,
        7909560,
        4617720,
        crop_top=crop_top,
        crop_bottom=0.035,
    )
    add_side_panel(slide, heading, bullets, conclusion)
    add_footer(slide, conclusion)
    return slide


def move_new_slides(prs, insert_after_index_1based, new_count):
    sld_id_lst = prs.slides._sldIdLst
    new_ids = list(sld_id_lst)[-new_count:]
    for sld_id in new_ids:
        sld_id_lst.remove(sld_id)
    insert_at = insert_after_index_1based
    for offset, sld_id in enumerate(new_ids):
        sld_id_lst.insert(insert_at + offset, sld_id)


def update_page_numbers(prs):
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if not text.isdigit():
                continue
            if shape.left > 11_000_000 and shape.top < 450_000:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = str(idx)
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(12)
                    run.font.bold = True
                    run.font.color.rgb = NAVY


def main():
    ppt = latest_target_ppt()
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup = ppt.with_name(f"{ppt.stem}.before_layered_map_insert_{timestamp}{ppt.suffix}")
    shutil.copy2(ppt, backup)
    print(f"backup: {backup}")

    prs = Presentation(str(ppt))
    logo_blob = None
    if len(prs.slides) >= 14:
        for shape in prs.slides[13].shapes:
            if shape.shape_type == 13 and shape.left > 9_000_000 and shape.top < 500_000:
                logo_blob = shape.image.blob
                break

    specs = [
        (
            "3.2A",
            "多源空间证据链：分层读图逻辑",
            "fig_total_layer_06_four_panel_evidence_chain.png",
            "讲述顺序",
            [
                "1. 空间骨架：区界、建筑/AOI、道路等级",
                "2. 服务供给：POI密度与设施类型",
                "3. 路径观测：步行路由与街景采样",
                "4. 风险解释：社区SAII与人口规模",
            ],
            "先分层后综合，形成从数据到指标的证据链。",
            15,
            0.13,
        ),
        (
            "3.2B",
            "空间基础层：路网等级与建筑/AOI",
            "fig_total_layer_01_base_road_hierarchy.png",
            "本页说明",
            [
                "• 主干/快速路决定跨片区通达骨架",
                "• 居住/服务道路决定街区内部连接",
                "• 步行/台阶/绿道构成慢行可达性基础",
                "• 建筑/AOI用于解释空间形态和阻隔",
            ],
            "这一层回答：居民能否沿连续路网接近服务。",
            16,
            0.13,
        ),
        (
            "3.2C",
            "服务供给层：POI密度与设施类型结构",
            "fig_total_layer_02_poi_service_supply.png",
            "本页说明",
            [
                "• 热度越高，代表服务供给越集中",
                "• 六类POI区分医疗、教育、交通、购物、餐饮、公共生活",
                "• 高密度不必然等于真实可达",
                "• 需与路网和夜间可用性共同解释",
            ],
            "这一层回答：服务在哪里，以及供给是否空间集聚。",
            17,
            0.13,
        ),
        (
            "3.2D",
            "路径证据层：步行路由与街景采样轨迹",
            "fig_total_layer_03_walk_route_streetview_sampling.png",
            "本页说明",
            [
                "• 青色线为可步行路由网络",
                "• 橙色点为20m间隔街景采样",
                "• 路由用于Network Ratio和M2SFCA",
                "• 街景用于GTA和步行环境质量校准",
            ],
            "这一层回答：服务看似近时，真实路径如何到达。",
            18,
            0.13,
        ),
        (
            "3.2E",
            "风险解释层：社区/AOI与SAII时间贫困",
            "fig_total_layer_04_community_aoi_saii_risk.png",
            "本页说明",
            [
                "• 圆点大小表示社区人口规模",
                "• 颜色越红，SAII时间贫困风险越高",
                "• 黑圈标记高风险社区",
                "• 高POI密度与高SAII并存，说明“近而难达”",
            ],
            "这一层回答：哪些社区承受了可达性幻觉。",
            19,
            0.13,
        ),
        (
            "3.2F",
            "低密度综合图：从图层证据进入指标计算",
            "fig_total_layer_05_low_density_synthesis.png",
            "过渡逻辑",
            [
                "• 只保留道路骨架、POI热区、路径轨迹和高风险社区",
                "• 避免再次堆叠全部点线面",
                "• 为后续Network Ratio、SAI、GTA、AII提供空间解释",
                "• 下一页进入网络距离与欧氏距离修正",
            ],
            "从这里开始，图层证据转化为可计算指标。",
            20,
            0.13,
        ),
    ]

    for spec in specs:
        add_map_slide(prs, logo_blob, *spec)

    move_new_slides(prs, insert_after_index_1based=14, new_count=len(specs))
    update_page_numbers(prs)
    prs.save(str(ppt))
    print(f"updated: {ppt}")
    print(f"inserted slides: {len(specs)} after slide 14")


if __name__ == "__main__":
    main()
