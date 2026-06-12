from __future__ import annotations

import shutil
from datetime import datetime
from pathlib import Path
from typing import NamedTuple

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
from matplotlib import font_manager
from matplotlib.lines import Line2D
from matplotlib.patches import Patch
from pptx import Presentation


ROOT = Path(__file__).resolve().parent
PPTX_PATH = ROOT / "哈工大PPT_博士答辩模板优化版_专业排版v8_跨区对比版(1).pptx"
WORK_DIR = ROOT / "_ppt_chart_work"
CHART_PATH = WORK_DIR / "obstacle_score_redrawn.png"
EXPORT_PATH = WORK_DIR / "slide33_shape14_redrawn_export.png"
FALLBACK_PPTX_PATH = PPTX_PATH.with_name(PPTX_PATH.stem + ".obstacle_legend_fixed.pptx")

DISTRICTS = ["南山区", "宝安区", "福田区", "龙华区"]
SCORES = np.array([8.45, 8.26, 3.62, 2.86])
COUNTS = [136, 58, 48, 48]
BOTTOM_BLOCK = np.array([11.8, 11.7, 4.8, 3.8])
BAR_COLOR = "#176d86"
LINE_COLOR = "#ef8a17"


class ReplaceResult(NamedTuple):
    backup_path: Path
    output_path: Path
    replaced_original: bool


def configure_fonts() -> None:
    candidates = [
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
        Path("C:/Windows/Fonts/simsun.ttc"),
    ]
    for font_path in candidates:
        if font_path.exists():
            font_manager.fontManager.addfont(str(font_path))
            if font_path.name.lower().startswith("msyh"):
                plt.rcParams["font.family"] = "Microsoft YaHei"
            elif font_path.name.lower().startswith("simhei"):
                plt.rcParams["font.family"] = "SimHei"
            else:
                plt.rcParams["font.family"] = "SimSun"
            break
    plt.rcParams["axes.unicode_minus"] = False


def draw_chart() -> None:
    configure_fonts()
    WORK_DIR.mkdir(exist_ok=True)

    x = np.arange(len(DISTRICTS))
    fig, ax = plt.subplots(figsize=(13.4, 8.08), dpi=220)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    bars = ax.bar(x, SCORES, width=0.62, color=BAR_COLOR, label="视觉障碍评分（柱）", zorder=3)
    for bar, score, count in zip(bars, SCORES, COUNTS):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            max(score - 0.22, 0.25),
            f"{score:.2f}\nn={count}",
            ha="center",
            va="top",
            fontsize=12.5,
            color="white",
            fontweight="bold",
            linespacing=1.1,
        )

    ax2 = ax.twinx()
    line = ax2.plot(
        x,
        BOTTOM_BLOCK,
        color=LINE_COLOR,
        linewidth=3.2,
        marker="o",
        markersize=8.5,
        label="底部遮挡占比（折线）",
        zorder=4,
    )[0]
    for i, value in enumerate(BOTTOM_BLOCK):
        offset = 0.18 if i < 2 else 0.22
        ax2.text(
            x[i],
            value + offset,
            f"{value:.1f}%",
            ha="center",
            va="bottom",
            fontsize=12.5,
            color="#b85f00",
            fontweight="bold",
        )

    ax.set_title(
        "跨区街景视觉障碍：南山/宝安显著高于福田/龙华",
        fontsize=21.5,
        color="#0f5f7a",
        pad=22,
        fontweight="bold",
    )
    ax.set_ylabel("视觉障碍评分（越高越强）", fontsize=15.5, color="#0f3044", labelpad=14)
    ax2.set_ylabel("底部遮挡占比（%）", fontsize=14.5, color=LINE_COLOR, labelpad=14)

    ax.set_xticks(x)
    ax.set_xticklabels(DISTRICTS, fontsize=17)
    ax.set_ylim(0, 9.2)
    ax.set_yticks(np.arange(0, 9.5, 1))
    ax2.set_ylim(3.4, 12.4)
    ax2.set_yticks(np.arange(4, 13, 1))

    ax.tick_params(axis="y", labelsize=12)
    ax2.tick_params(axis="y", labelsize=12, colors="#111111")
    ax.grid(axis="y", linestyle="--", linewidth=0.8, alpha=0.35, zorder=0)
    ax.spines["top"].set_visible(False)
    ax2.spines["top"].set_visible(False)

    handles = [
        Patch(facecolor=BAR_COLOR, edgecolor="none", label="视觉障碍评分（柱）"),
        Line2D([0], [0], color=LINE_COLOR, marker="o", linewidth=3.2, markersize=8, label="底部遮挡占比（折线）"),
    ]
    ax.legend(
        handles=handles,
        loc="upper right",
        frameon=False,
        fontsize=13,
        bbox_to_anchor=(0.985, 0.94),
        handlelength=2.0,
    )

    fig.tight_layout(rect=(0.02, 0.035, 0.98, 0.98))
    fig.savefig(CHART_PATH, dpi=220, facecolor="white")
    plt.close(fig)


def replace_slide_picture() -> ReplaceResult:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = PPTX_PATH.with_name(PPTX_PATH.stem + f".before_obstacle_chart_redraw_{timestamp}.pptx")
    shutil.copy2(PPTX_PATH, backup_path)

    prs = Presentation(PPTX_PATH)
    slide = prs.slides[32]
    target = None
    for shape in slide.shapes:
        if shape.shape_type == 13 and shape.left == 768096:
            target = shape
            break
        if shape.shape_id == 14 or shape.name == "Picture 13":
            target = shape
            break
    if target is None:
        raise RuntimeError("Could not find slide 33 obstacle chart picture")

    left, top, width, height = target.left, target.top, target.width, target.height
    element = target._element
    element.getparent().remove(element)
    slide.shapes.add_picture(str(CHART_PATH), left, top, width=width, height=height)
    temp_path = PPTX_PATH.with_name(PPTX_PATH.stem + f".obstacle_legend_fixed_tmp_{timestamp}.pptx")
    prs.save(temp_path)
    try:
        shutil.move(str(temp_path), str(PPTX_PATH))
        return ReplaceResult(backup_path, PPTX_PATH, True)
    except PermissionError:
        if FALLBACK_PPTX_PATH.exists():
            FALLBACK_PPTX_PATH.unlink()
        shutil.move(str(temp_path), str(FALLBACK_PPTX_PATH))
        return ReplaceResult(backup_path, FALLBACK_PPTX_PATH, False)


def export_replacement_for_check(pptx_path: Path) -> None:
    prs = Presentation(pptx_path)
    slide = prs.slides[32]
    for shape in slide.shapes:
        if shape.shape_type == 13 and shape.left == 768096:
            EXPORT_PATH.write_bytes(shape.image.blob)
            return
    raise RuntimeError("Could not export replaced obstacle chart picture")


def main() -> None:
    draw_chart()
    result = replace_slide_picture()
    export_replacement_for_check(result.output_path)
    print(f"chart={CHART_PATH}")
    print(f"backup={result.backup_path}")
    print(f"output={result.output_path}")
    print(f"replaced_original={result.replaced_original}")
    print(f"export={EXPORT_PATH}")


if __name__ == "__main__":
    main()
