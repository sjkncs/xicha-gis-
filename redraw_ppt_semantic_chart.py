from __future__ import annotations

import json
import re
import shutil
from datetime import datetime
from pathlib import Path

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib import font_manager
from pptx import Presentation


ROOT = Path(__file__).resolve().parent
PPTX_PATH = ROOT / "哈工大PPT_博士答辩模板优化版_专业排版v8_跨区对比版(1).pptx"
SEG_PATH = ROOT / "streetview-analysis" / "gpu_scripts" / "results" / "seg_results.csv"
YOLO_PATH = ROOT / "自选年份" / "yolo_results_merged.json"
WORK_DIR = ROOT / "_ppt_chart_work"
CHART_PATH = WORK_DIR / "semantic_structure_redrawn.png"
EXPORT_PATH = WORK_DIR / "slide33_shape19_redrawn_export.png"

DISTRICTS = ["南山区", "宝安区", "福田区", "龙华区"]
SERIES = [
    ("建筑界面", "pct_building", "#176d86"),
    ("道路", "pct_road", "#bd303b"),
    ("绿视率", "pct_green", "#459875"),
    ("天空开敞", "pct_sky", "#f1992d"),
]


def configure_fonts() -> None:
    candidates = [
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("C:/Windows/Fonts/simhei.ttf"),
        Path("C:/Windows/Fonts/simsun.ttc"),
    ]
    for font_path in candidates:
        if font_path.exists():
            font_manager.fontManager.addfont(str(font_path))
            if font_path.name.lower().startswith("msyh"):
                plt.rcParams["font.family"] = "Microsoft YaHei"
            elif font_path.name.lower().startswith("simhei"):
                plt.rcParams["font.family"] = "SimHei"
            else:
                plt.rcParams["font.family"] = "SimSun"
            break
    plt.rcParams["axes.unicode_minus"] = False


def load_semantic_matrix() -> pd.DataFrame:
    seg = pd.read_csv(SEG_PATH)
    yolo = json.loads(YOLO_PATH.read_text(encoding="utf-8"))
    district_re = re.compile(r"(南山区|宝安区|福田区|龙华区|Village)")

    mapped_rows = []
    for row in yolo:
        match = district_re.search(row.get("file", ""))
        if match:
            mapped_rows.append({"pano_name": row.get("filename"), "district": match.group(1)})

    district_map = pd.DataFrame(mapped_rows).drop_duplicates("pano_name")
    merged = seg.merge(district_map, on="pano_name", how="left")
    if merged["district"].isna().any():
        missing = int(merged["district"].isna().sum())
        raise RuntimeError(f"{missing} segmentation rows could not be mapped to districts")

    cols = [col for _, col, _ in SERIES]
    matrix = merged.groupby("district")[cols].mean().loc[DISTRICTS]
    return matrix


def draw_chart(matrix: pd.DataFrame) -> None:
    configure_fonts()
    WORK_DIR.mkdir(exist_ok=True)

    x = np.arange(len(DISTRICTS))
    width = 0.18
    offsets = np.linspace(-1.5 * width, 1.5 * width, len(SERIES))

    fig, ax = plt.subplots(figsize=(13.4, 8.05), dpi=220)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    max_val = 0.0
    for offset, (label, col, color) in zip(offsets, SERIES):
        values = matrix[col].to_numpy()
        max_val = max(max_val, float(values.max()))
        bars = ax.bar(x + offset, values, width=width, label=label, color=color)
        for bar, value in zip(bars, values):
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                bar.get_height() + 0.35,
                f"{value:.1f}",
                ha="center",
                va="bottom",
                fontsize=11.5,
                color="#111111",
            )

    ax.set_title(
        "跨区街景语义结构：开放度、绿量与建筑界面的差异",
        fontsize=22,
        color="#0f5f7a",
        pad=22,
        fontweight="regular",
    )
    ax.set_ylabel("DeepLabV3语义分割占比（%）", fontsize=15, color="#0f3044", labelpad=14)
    ax.set_xticks(x)
    ax.set_xticklabels(DISTRICTS, fontsize=17)
    ax.tick_params(axis="y", labelsize=12)
    ax.set_ylim(0, max_val + 7)
    ax.set_yticks(np.arange(0, 46, 5))
    ax.grid(axis="y", linestyle="--", linewidth=0.8, alpha=0.35)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    legend = ax.legend(
        ncol=4,
        loc="upper center",
        bbox_to_anchor=(0.5, -0.095),
        frameon=False,
        fontsize=13,
        columnspacing=2.3,
        handlelength=1.7,
    )
    for text in legend.get_texts():
        text.set_color("#111111")

    fig.tight_layout(rect=(0.02, 0.055, 0.99, 0.98))
    fig.savefig(CHART_PATH, dpi=220, facecolor="white")
    plt.close(fig)


def replace_slide_picture() -> Path:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_path = PPTX_PATH.with_name(PPTX_PATH.stem + f".before_semantic_chart_redraw_{timestamp}.pptx")
    shutil.copy2(PPTX_PATH, backup_path)

    prs = Presentation(PPTX_PATH)
    slide = prs.slides[32]
    target = None
    for shape in slide.shapes:
        if shape.shape_id == 19 or shape.name == "Picture 18":
            target = shape
            break
    if target is None:
        raise RuntimeError("Could not find slide 33 semantic chart picture")

    left, top, width, height = target.left, target.top, target.width, target.height
    element = target._element
    element.getparent().remove(element)
    slide.shapes.add_picture(str(CHART_PATH), left, top, width=width, height=height)
    prs.save(PPTX_PATH)
    return backup_path


def export_replacement_for_check() -> None:
    prs = Presentation(PPTX_PATH)
    slide = prs.slides[32]
    candidates = []
    for shape in slide.shapes:
        if getattr(shape, "image", None) is None:
            continue
        try:
            blob = shape.image.blob
        except Exception:
            continue
        candidates.append((len(blob), blob))
    if not candidates:
        raise RuntimeError("No pictures found on slide 33 after replacement")
    blob = max(candidates, key=lambda item: item[0])[1]
    EXPORT_PATH.write_bytes(blob)


def main() -> None:
    matrix = load_semantic_matrix()
    print(matrix.round(3).to_string())
    draw_chart(matrix)
    backup_path = replace_slide_picture()
    export_replacement_for_check()
    print(f"chart={CHART_PATH}")
    print(f"backup={backup_path}")
    print(f"export={EXPORT_PATH}")


if __name__ == "__main__":
    main()
